"""
extrair_metadados.py
======================
PASSO 1 do pipeline de biblioteca.

Lê todos os documentos de PASTA_ORIGEM (recursivo, config.py) e cria/atualiza
um catálogo central `catalogo.json` com, por documento: titulo, autor, resumo,
tag_funcao e keywords multilingue.

Dois grupos de ficheiros (config.py):
 - EXTENSOES_SUPORTADAS: texto extraído localmente (pdf, docx, pptx, xlsx, txt...)
 - EXTENSOES_NOME_APENAS: sem texto extraível (rvt, dwg, jpg...) — classificados
   pelo nome do ficheiro + pasta de origem

É RESUMÍVEL: ficheiros já indexados (mesmo hash) são ignorados. Corre outra
vez sempre que adicionares livros novos a PASTA_ORIGEM.

Instalação:
 pip install openai python-docx python-pptx pypdf openpyxl
"""

import hashlib
import json
from pathlib import Path
from datetime import datetime

# MIGRAÇÃO: google-genai -> OpenRouter via llm_util
from llm_util import pedir_json

from config import (
    PASTA_ORIGEM, CATALOGO_PATH,
    LINGUAS_KEYWORDS, TAGS_FUNCAO, PAGINAS_PDF, TAMANHO_LOTE,
    EXTENSOES_TODAS,
)

# ---------------------------------------------------------------------------
# HASH
# ---------------------------------------------------------------------------

def hash_ficheiro(caminho: Path, bloco: int = 8192) -> str:
    h = hashlib.sha256()
    with open(caminho, "rb") as f:
        while chunk := f.read(bloco):
            h.update(chunk)
    return h.hexdigest()

# ---------------------------------------------------------------------------
# NORMALIZAÇÃO DE KEYWORDS (defesa contra desvios do modelo)
# ---------------------------------------------------------------------------

def achatar_termos(t) -> list:
    """Normaliza o valor de keywords de uma língua para lista plana de strings.
    Apanha listas aninhadas, strings soltas e lixo — o catálogo nunca guarda malformado."""
    if isinstance(t, str):
        return [t.strip()] if t.strip() else []
    if isinstance(t, list):
        termos = []
        for item in t:
            termos.extend(achatar_termos(item))
        return termos
    return []

def keywords_validas(kw_raw) -> dict:
    """Garante o formato {lingua: [strings]} a partir do que o modelo devolver."""
    if not isinstance(kw_raw, dict):
        plano = achatar_termos(kw_raw)
        return {"pt": plano} if plano else {l: [] for l in LINGUAS_KEYWORDS}
    resultado = {}
    for l, t in kw_raw.items():
        resultado[str(l)] = achatar_termos(t)
    return resultado

# ---------------------------------------------------------------------------
# EXTRAÇÃO DE TEXTO POR TIPO DE FICHEIRO
# ---------------------------------------------------------------------------

def extrair_texto_pdf(caminho: Path, n_paginas: int = PAGINAS_PDF) -> str:
    from pypdf import PdfReader
    try:
        leitor = PdfReader(str(caminho))
    except Exception:
        return ""
    partes = []
    for pagina in leitor.pages[:n_paginas]:
        try:
            partes.append(pagina.extract_text() or "")
        except Exception:
            continue
    return "\n".join(partes)

def extrair_texto_docx(caminho: Path) -> str:
    from docx import Document
    doc = Document(caminho)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extrair_texto_pptx(caminho: Path) -> str:
    from pptx import Presentation
    prs = Presentation(caminho)
    partes = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                partes.append(shape.text_frame.text)
    return "\n".join(partes)

def extrair_texto_xlsx(caminho: Path) -> str:
    from openpyxl import load_workbook
    try:
        wb = load_workbook(caminho, read_only=True, data_only=True)
    except Exception:
        return ""
    partes = []
    for ws in wb.worksheets[:3]:                      # primeiras 3 folhas chegam
        for row in ws.iter_rows(max_row=50, values_only=True):
            linha = " ".join(str(c) for c in row if c is not None)
            if linha.strip():
                partes.append(linha)
    return "\n".join(partes)[:4000]

def extrair_texto(caminho: Path) -> str:
    """Nunca deixa uma exceção sair daqui — um ficheiro problemático (corrompido,
    ou com extensão errada, ex: .doc antigo renomeado para .docx) fica só com
    texto vazio e regista um aviso, em vez de parar o pipeline inteiro.
    Ficheiros sem extrator (CAD/BIM, imagens) devolvem "" por desenho — são
    classificados pelo nome + pasta de origem."""
    ext = caminho.suffix.lower()
    try:
        if ext == ".pdf":
            return extrair_texto_pdf(caminho)
        if ext == ".docx":
            return extrair_texto_docx(caminho)
        if ext == ".pptx":
            return extrair_texto_pptx(caminho)
        if ext == ".xlsx":
            return extrair_texto_xlsx(caminho)
        if ext in {".txt", ".md", ".rtf"}:
            return caminho.read_text(encoding="utf-8", errors="ignore")[:4000]
    except Exception as e:
        print(f"⚠️ Não consegui ler o conteúdo de '{caminho.name}' ({type(e).__name__}: {e}). "
              f"Verifica se o ficheiro não está corrompido ou com a extensão errada "
              f"(ex: um .doc antigo renomeado para .docx). A continuar sem texto deste ficheiro.")
    return ""

# ---------------------------------------------------------------------------
# CATÁLOGO
# ---------------------------------------------------------------------------

def carregar_catalogo() -> dict:
    if CATALOGO_PATH.exists():
        with open(CATALOGO_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def guardar_catalogo(catalogo: dict):
    CATALOGO_PATH.parent.mkdir(parents=True, exist_ok=True)
    with open(CATALOGO_PATH, "w", encoding="utf-8") as f:
        json.dump(catalogo, f, ensure_ascii=False, indent=2)

# ---------------------------------------------------------------------------
# CLASSIFICAÇÃO EM LOTE (LLM via OpenRouter)
# ---------------------------------------------------------------------------

def construir_prompt_lote(linguas: list) -> str:
    linguas_str = ", ".join(linguas)
    opcoes_funcao = ", ".join(TAGS_FUNCAO)
    return f"""Para cada documento apresentado abaixo (título do ficheiro + pasta de
origem atual, como pista + excerto de texto), extrai metadados bibliográficos.
Responde APENAS com uma lista JSON válida, um objeto por documento, na MESMA
ORDEM em que foram apresentados:

[
 {{
   "titulo": "título real do documento (não o nome do ficheiro, a menos que coincida)",
   "autor": "autor(es) ou 'desconhecido'",
   "resumo": "resumo objetivo em português, 2-3 frases",
   "tag_funcao": "uma destas opções: {opcoes_funcao}",
   "keywords": {{
     {", ".join(f'"{l}": ["...", "..."]' for l in linguas)}
   }}
 }},
 ...
]

"tag_funcao" descreve O QUE o documento É (não o tópico técnico de que fala) —
usa a pasta de origem como pista mas confirma com o conteúdo (ex: um ficheiro
na pasta "D-WEBINARS" é provavelmente "Webinar"; um na pasta "B-EXAMPLES" é
provavelmente "Exemplo"; um livro técnico normal é "Livro/Manual").

Para documentos SEM excerto de texto (ficheiros CAD/BIM como .rvt/.dwg, ou
imagens), classifica pelo nome do ficheiro e pasta de origem; se o nome não
for informativo (ex: IMG_2047.jpg), usa tag_funcao "Outro" e keywords genéricas.

Cada lista de keywords deve ter entre 5 a 10 termos técnicos relevantes,
traduzidos/equivalentes nas línguas: {linguas_str}. Não incluas texto fora do JSON.
"""

def classificar_lote(documentos: list) -> list:
    """documentos: lista de (caminho, texto_extraido, pasta_origem)."""
    prompt = construir_prompt_lote(LINGUAS_KEYWORDS)
    for caminho, texto, pasta_origem in documentos:
        prompt += (f"\n\n--- Documento: {caminho.name} "
                   f"(pasta de origem: {pasta_origem}) ---\n{texto[:4000]}")

    bruto = pedir_json(prompt)  # MIGRAÇÃO: já devolve texto limpo, sem fences
    try:
        parsed = json.loads(bruto)
        # NemoTron às vezes devolve objeto único ou string em vez de lista
        if isinstance(parsed, dict):
            parsed = [parsed]
        if not isinstance(parsed, list) or not all(isinstance(x, dict) for x in parsed):
            raise ValueError(f"JSON não é lista de objetos: {type(parsed).__name__}")
        return parsed
    except (json.JSONDecodeError, ValueError) as e:
        print(f"⚠️ Falha a interpretar resposta do modelo neste lote ({e}) — a saltar.")
        print("--- Resposta em bruto (para diagnóstico) ---")
        print(bruto[:1000])
        print("--- fim da resposta em bruto ---")
        return [{"titulo": c.stem, "autor": "desconhecido", "resumo": "",
                 "tag_funcao": "Outro",
                 "keywords": {l: [] for l in LINGUAS_KEYWORDS}} for c, _, _ in documentos]

# ---------------------------------------------------------------------------
# ORQUESTRAÇÃO
# ---------------------------------------------------------------------------

def main():
    catalogo = carregar_catalogo()
    hashes_existentes = {v["hash"] for v in catalogo.values()}

    todos = [f for f in PASTA_ORIGEM.rglob("*") if f.suffix.lower() in EXTENSOES_TODAS]
    print(f"Encontrados {len(todos)} documentos em {PASTA_ORIGEM}.")

    por_indexar = []
    for f in todos:
        h = hash_ficheiro(f)
        if h in hashes_existentes:
            continue
        por_indexar.append((f, h))

    print(f"Já indexados: {len(todos) - len(por_indexar)} | Novos: {len(por_indexar)}\n")

    for i in range(0, len(por_indexar), TAMANHO_LOTE):
        lote = por_indexar[i:i + TAMANHO_LOTE]
        print(f"A processar lote {i // TAMANHO_LOTE + 1} ({len(lote)} documentos)...")

        documentos_com_texto = []
        for caminho, _ in lote:
            texto = extrair_texto(caminho)
            pasta_origem = caminho.parent.name
            documentos_com_texto.append((caminho, texto, pasta_origem))

        metadados_lote = classificar_lote(documentos_com_texto)
        # NemoTron e outros nem sempre devolvem lista de dicts -> normalizar
        if isinstance(metadados_lote, dict):
            metadados_lote = [metadados_lote]
        if not isinstance(metadados_lote, list):
            metadados_lote = list(metadados_lote) if isinstance(metadados_lote, (tuple, set)) else []

        for (caminho, h), metadados in zip(lote, metadados_lote):
            if not isinstance(metadados, dict):
                print(f"⚠️ Metadados inválidos para '{caminho.name}' ({type(metadados).__name__}: {str(metadados)[:200]}). A usar defaults.")
                metadados = {}
            chave = str(caminho.relative_to(PASTA_ORIGEM))
            catalogo[chave] = {
                "caminho": str(caminho),
                "hash": h,
                "tipo_ficheiro": caminho.suffix.lower().lstrip("."),
                "pasta_origem_atual": caminho.parent.name,
                "titulo": metadados.get("titulo", caminho.stem),
                "autor": metadados.get("autor", "desconhecido"),
                "resumo": metadados.get("resumo", ""),
                "tag_funcao": metadados.get("tag_funcao", "Outro"),
                "keywords": keywords_validas(metadados.get("keywords", {})),  # nunca guarda malformado
                "data_indexado": datetime.now().isoformat(timespec="seconds"),
            }

        guardar_catalogo(catalogo)  # progressivo — nada se perde se algo falhar a meio

    print(f"\nConcluído. Catálogo tem {len(catalogo)} documentos em {CATALOGO_PATH}")
    return True

if __name__ == "__main__":
    from llm_util import LimiteDiarioAtingido  # MIGRAÇÃO
    try:
        main()
    except LimiteDiarioAtingido as e:
        print(f"\n🛑 {e}")