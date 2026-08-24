"""
adicionar_pasta.py
====================
SCRIPT À PARTE — não faz parte da sequência do correr_tudo.py. Corres isto
manualmente sempre que tiveres uma pasta nova de conteúdo para juntar à
biblioteca já organizada (config.PASTA_DESTINO), sem duplicar o que já lá
está. Ajusta config.ORIGEM_NOVA antes de correr.

Deteta duplicados por hash contra tudo o que já existe no destino, e para os
ficheiros genuinamente novos: extrai metadados, escolhe pasta de destino de
entre as que já existem (só propõe pasta nova se não encaixar em nenhuma),
renomeia já com tags, e move.
"""

import hashlib
import json
import re
import shutil
from pathlib import Path
from datetime import datetime

# MIGRAÇÃO: google-genai -> OpenRouter via llm_util
from llm_util import pedir_json

from config import (  # MIGRAÇÃO: era 'from config_gemini import' — bug corrigido
    ORIGEM_NOVA, PASTA_DESTINO, CATALOGO_PATH, INDICE_HASHES_PATH,
    DRY_RUN, TAMANHO_LOTE,
    PROFUNDIDADE_TAXONOMIA, LINGUAS_KEYWORDS, TAGS_FUNCAO,
    MAX_TAGS_POR_FICHEIRO, EXTENSOES_SUPORTADAS,
)

# ---------------------------------------------------------------------------
# HASH / ÍNDICE (cacheado — só faz scan completo do destino uma vez)
# ---------------------------------------------------------------------------

def hash_ficheiro(caminho: Path, bloco: int = 8192) -> str:
    h = hashlib.sha256()
    with open(caminho, "rb") as f:
        while chunk := f.read(bloco):
            h.update(chunk)
    return h.hexdigest()

def construir_ou_carregar_indice() -> dict:
    if INDICE_HASHES_PATH.exists():
        with open(INDICE_HASHES_PATH, "r", encoding="utf-8") as f:
            print(f"Índice de hashes carregado de cache ({INDICE_HASHES_PATH}).")
            return json.load(f)

    print("Sem cache — a construir índice de hashes do destino (só acontece uma vez)...")
    indice = {}
    todos = list(PASTA_DESTINO.rglob("*"))
    for i, f in enumerate(todos):
        if f.is_file():
            try:
                indice[hash_ficheiro(f)] = str(f)
            except (OSError, PermissionError):
                continue
        if i % 500 == 0:
            print(f"  ...{i}/{len(todos)} processados")

    guardar_indice(indice)
    print(f"Índice construído com {len(indice)} ficheiros.")
    return indice

def guardar_indice(indice: dict):
    INDICE_HASHES_PATH.parent.mkdir(parents=True, exist_ok=True)
    with open(INDICE_HASHES_PATH, "w", encoding="utf-8") as f:
        json.dump(indice, f, ensure_ascii=False, indent=2)

# ---------------------------------------------------------------------------
# TAXONOMIA EXISTENTE (lida diretamente da árvore de pastas do destino)
# ---------------------------------------------------------------------------

def listar_taxonomia_existente(profundidade: int = PROFUNDIDADE_TAXONOMIA) -> list:
    pastas = []

    def _percorrer(pasta: Path, nivel: int, prefixo: str):
        if nivel > profundidade:
            return
        for item in sorted(pasta.iterdir()):
            if item.is_dir() and not item.name.startswith("_"):
                caminho_rel = f"{prefixo}{item.name}"
                pastas.append(caminho_rel)
                _percorrer(item, nivel + 1, caminho_rel + "/")

    _percorrer(PASTA_DESTINO, 1, "")
    return pastas

# ---------------------------------------------------------------------------
# EXTRAÇÃO DE TEXTO
# ---------------------------------------------------------------------------

def extrair_texto_pdf(caminho: Path, n_paginas: int = 12) -> str:
    from pypdf import PdfReader
    try:
        leitor = PdfReader(str(caminho))
    except Exception:
        return ""
    partes = []
    for pagina in leitor.pages[:n_paginas]:
        try:
            partes.append(pagina.extract_text() or "")
        except Exception:
            continue
    return "\n".join(partes)

def extrair_texto_docx(caminho: Path) -> str:
    from docx import Document
    doc = Document(caminho)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extrair_texto_pptx(caminho: Path) -> str:
    from pptx import Presentation
    prs = Presentation(caminho)
    partes = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                partes.append(shape.text_frame.text)
    return "\n".join(partes)

def extrair_texto(caminho: Path) -> str:
    ext = caminho.suffix.lower()
    if ext == ".pdf":
        return extrair_texto_pdf(caminho)
    if ext == ".docx":
        return extrair_texto_docx(caminho)
    if ext == ".pptx":
        return extrair_texto_pptx(caminho)
    return ""

# ---------------------------------------------------------------------------
# CLASSIFICAÇÃO EM LOTE
# ---------------------------------------------------------------------------

def construir_prompt(taxonomia: list) -> str:
    linguas_str = ", ".join(LINGUAS_KEYWORDS)
    opcoes_funcao = ", ".join(TAGS_FUNCAO)
    lista_pastas = "\n".join(f" - {p}" for p in taxonomia)
    return f"""Já existe uma biblioteca organizada com estas pastas:

{lista_pastas}

Para cada documento novo abaixo, extrai metadados E escolhe a pasta de
destino mais adequada de ENTRE AS QUE JÁ EXISTEM acima. Só propõe pasta nova
se genuinamente não encaixar em nenhuma — nesse caso, prefixa com "NOVA: ".

Responde APENAS com uma lista JSON, um objeto por documento, na MESMA ORDEM:
[
  {{
    "titulo": "...", "autor": "... ou desconhecido",
    "resumo": "2-3 frases em português",
    "tag_funcao": "uma destas: {opcoes_funcao}",
    "keywords": {{{", ".join(f'"{l}": ["...", "..."]' for l in LINGUAS_KEYWORDS)}}},
    "pasta_destino": "caminho exato de uma pasta listada, ou NOVA: ..."
  }},
  ...
]
Keywords: 5-10 termos técnicos por língua ({linguas_str}). Não incluas texto fora do JSON.
"""

def classificar_lote(documentos: list, taxonomia: list) -> list:
    prompt = construir_prompt(taxonomia)
    for caminho, texto in documentos:
        prompt += f"\n\n--- Documento: {caminho.name} ---\n{texto[:4000]}"

    bruto = pedir_json(prompt)  # MIGRAÇÃO
    try:
        return json.loads(bruto)
    except json.JSONDecodeError:
        print("⚠️ Falha a interpretar resposta do modelo neste lote — a saltar.")
        print("--- Resposta em bruto (para diagnóstico) ---")
        print(bruto[:1000])
        print("--- fim da resposta em bruto ---")
        return [{"titulo": c.stem, "autor": "desconhecido", "resumo": "",
                 "tag_funcao": "Outro", "keywords": {}, "pasta_destino": "NOVA: Outros"}
                for c, _ in documentos]

# ---------------------------------------------------------------------------
# TAGS NO NOME (convenção TagSpaces)
# ---------------------------------------------------------------------------

def sanitizar_tag(tag: str) -> str:
    tag = tag.strip().lower()
    tag = re.sub(r"\s+", "-", tag)
    tag = re.sub(r'[\\/:*?"<>|\[\]]', "", tag)
    return tag

def escolher_tags(dados: dict, max_tags: int = MAX_TAGS_POR_FICHEIRO) -> list:
    tags = []
    tag_funcao = dados.get("tag_funcao")
    if tag_funcao and tag_funcao != "Outro":
        tags.append(sanitizar_tag(tag_funcao))
    for kw in dados.get("keywords", {}).get("pt", []):
        if len(tags) >= max_tags:
            break
        tag = sanitizar_tag(kw)
        if tag and tag not in tags:
            tags.append(tag)
    return tags

def construir_novo_nome(caminho: Path, tags: list) -> str:
    if not tags:
        return caminho.name
    return f"{caminho.stem} [{' '.join(tags)}]{caminho.suffix}"

# ---------------------------------------------------------------------------
# ORQUESTRAÇÃO
# ---------------------------------------------------------------------------

def carregar_catalogo() -> dict:
    if CATALOGO_PATH.exists():
        with open(CATALOGO_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def guardar_catalogo(catalogo: dict):
    CATALOGO_PATH.parent.mkdir(parents=True, exist_ok=True)
    with open(CATALOGO_PATH, "w", encoding="utf-8") as f:
        json.dump(catalogo, f, ensure_ascii=False, indent=2)

def main() -> bool:
    indice = construir_ou_carregar_indice()
    catalogo = carregar_catalogo()
    taxonomia = listar_taxonomia_existente()
    print(f"Taxonomia existente no destino: {len(taxonomia)} pastas.\n")

    candidatos = [f for f in ORIGEM_NOVA.rglob("*")
                  if f.is_file() and f.suffix.lower() in EXTENSOES_SUPORTADAS]
    print(f"Ficheiros novos encontrados na pasta raiz: {len(candidatos)}\n")

    duplicados, a_processar = [], []
    for f in candidatos:
        h = hash_ficheiro(f)
        if h in indice:
            duplicados.append((f, indice[h]))
        else:
            a_processar.append((f, h))

    print(f"Já existem na biblioteca (duplicados, NÃO copiados): {len(duplicados)}")
    for novo, existente in duplicados:
        print(f"  • {novo.name} já existe como {existente}")

    print(f"\nA processar {len(a_processar)} ficheiros genuinamente novos...\n")

    for i in range(0, len(a_processar), TAMANHO_LOTE):
        lote = a_processar[i:i + TAMANHO_LOTE]
        print(f"Lote {i // TAMANHO_LOTE + 1} ({len(lote)} ficheiros)...")

        documentos = [(caminho, extrair_texto(caminho)) for caminho, _ in lote]
        classificacoes = classificar_lote(documentos, taxonomia)

        for (caminho, h), dados in zip(lote, classificacoes):
            pasta_destino_str = dados.get("pasta_destino", "NOVA: Outros")
            nova = pasta_destino_str.startswith("NOVA:")
            pasta_rel = pasta_destino_str.replace("NOVA:", "").strip()

            tags = escolher_tags(dados)
            novo_nome = construir_novo_nome(caminho, tags)
            destino_final = PASTA_DESTINO / pasta_rel / novo_nome

            marcador = "🆕 pasta nova" if nova else "pasta existente"
            if DRY_RUN:
                print(f"  [SIMULAÇÃO] ({marcador}) {caminho.name}")
                print(f"      -> {destino_final}")
            else:
                destino_final.parent.mkdir(parents=True, exist_ok=True)
                shutil.move(str(caminho), str(destino_final))
                print(f"  movido ({marcador}): {novo_nome} -> {destino_final.parent}")

            chave = str(destino_final.relative_to(PASTA_DESTINO))
            catalogo[chave] = {
                "caminho": str(destino_final),
                "hash": h,
                "tipo_ficheiro": caminho.suffix.lower().lstrip("."),
                "titulo": dados.get("titulo", caminho.stem),
                "autor": dados.get("autor", "desconhecido"),
                "resumo": dados.get("resumo", ""),
                "tag_funcao": dados.get("tag_funcao", "Outro"),
                "keywords": dados.get("keywords", {}),
                "data_indexado": datetime.now().isoformat(timespec="seconds"),
            }
            indice[h] = str(destino_final)
            if nova:
                taxonomia.append(pasta_rel)

        if not DRY_RUN:
            guardar_catalogo(catalogo)
            guardar_indice(indice)

    print(f"\nConcluído. {len(duplicados)} duplicados ignorados, "
          f"{len(a_processar)} processados.",
          "(modo simulação — nada foi movido)" if DRY_RUN else "")
    return not DRY_RUN

if __name__ == "__main__":
    from llm_util import LimiteDiarioAtingido  # MIGRAÇÃO
    try:
        main()
    except LimiteDiarioAtingido as e:
        print(f"\n🛑 {e}")